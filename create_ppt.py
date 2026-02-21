from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, content_lines, is_title_slide=False):
    if is_title_slide:
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        subtitle_text = "\n".join(content_lines)
        subtitle_shape.text = subtitle_text
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            
    else:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        tf.clear() # Clear default paragraph
        for i, line in enumerate(content_lines):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            if line.startswith("- "):
                p.level = 1
                p.text = line[2:]
            elif line.startswith("  - "):
                p.level = 2
                p.text = line[4:]

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    add_slide(prs, "Good Software – Cloud Suite\n릴리스 노트", [
        "버전: 3.2.0 | 출시일: 2025년 7월 25일",
        "대상 독자: 최종 사용자, IT 관리자, DevOps 팀, 구현 파트너"
    ], is_title_slide=True)
    
    # About
    add_slide(prs, "이 릴리스에 대하여", [
        "버전 3.2.0은 플랫폼 성능, 사용자 인터페이스, 시스템 보안에 개선을 제공합니다.",
        "새로운 워크플로우 빌더, 스마트 대시보드 및 향상된 기능 도입.",
        "버전 3.1.x와 완전히 하위 호환됩니다."
    ])
    
    # Highlights
    add_slide(prs, "릴리스 주요 사항", [
        "- 새로운 사용자 정의 워크플로우 빌더 – 코드 없는 자동화 프로세스",
        "- 고급 SAML 2.0 SSO 지원 – 보완성 강화",
        "- 성능 최적화 – 응답시간 단축 및 UI 로드 시간 최소화",
        "- 주요 버그 수정 – 안정성 향상",
        "- 기능 종료 공지 – 기존 도구 지원 중단"
    ])
    
    # New Features
    add_slide(prs, "새로운 기능", [
        "사용자 정의 워크플로우 빌더:",
        "- 드래그 앤 드롭 UI를 이용해 코드 없이 일상 업무 자동화",
        "역할 기반 대시보드:",
        "- 업무 역할에 따른 개인화된 화면으로 통찰력 제공",
        "감사 로그 내보내기:",
        "- 감사 데이터를 CSV 혹은 JSON으로 내보내 규정 준수를 지원"
    ])

    # Enhancements
    add_slide(prs, "개선 사항", [
        "REST API v2:",
        "- 평균 응답 시간이 40% 감소, Pagination 적용 범위 확대",
        "UI/UX:",
        "- 접근성 개선(WCAG 2.1 준수) 및 일관된 인터페이스 렌더링",
        "보안:",
        "- OAuth 2.0 토큰 만료 시간 직접 관리/SAML 비활성 자동 종료"
    ])
    
    # Deprecated & Known Issues
    add_slide(prs, "지원 중단 안내 & 알려진 문제", [
        "지원 중단 예정 기능:",
        "- 기존 CSV 내보내기 도구, 레거시 관리자 UI",
        "알려진 문제:",
        "- iOS 로드 지연 현상 등은 v3.2.1 패치에서 수정 예정"
    ])
    
    # Upgrades
    add_slide(prs, "업그레이드 및 요구사항", [
        "시스템 요구사항:",
        "- Ubuntu 22.04, PostgreSQL 13+, DB/OS 최신버전 유지 권장",
        "업그레이드 절차:",
        "1. 데이터베이스 및 시스템 환경 백업",
        "2. 설치 프로그램 실행 (관리자 패널)",
        "3. 상태 확인 및 진단(예상시간: 약 10분)"
    ])
    
    prs.save("Good.Software.Release.Notes_ko.pptx")
    print("PPT created successfully.")

if __name__ == "__main__":
    create_presentation()
